"""Adaptive document parser (Phase 3).

Selects the optimal parser based on file extension and produces a
uniform ``parsed_doc`` shape:

    {
        "pages": [
            {"page_num": int, "text": str, "tables": [{"markdown": str, "caption": str}],
             "image_refs": [str]},
            ...
        ],
        "metadata": {"title": str, "num_pages": int, "file_type": str, "source_path": str},
        "text_markdown": str,  # full markdown export when available
    }

Parsers:
- PDF   -> Docling (table-aware) with PyMuPDF fallback
- DOCX  -> python-docx with paragraph and table extraction
- PPTX  -> python-pptx with per-slide text and tables
- Other -> plain-text fallback (for .txt, .md, .html)
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List

logger = logging.getLogger(__name__)


class AdaptiveDocumentParser:
    """Dispatch by file extension."""

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md", ".html"}

    def __init__(self) -> None:
        self._docling = None
        self._docling_available = self._init_docling()

    # ------------------------------------------------------------------
    # Public
    # ------------------------------------------------------------------

    def parse(self, file_path: str) -> Dict[str, Any]:
        ext = Path(file_path).suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file format: {ext}")

        if ext == ".pdf":
            return self._parse_pdf(file_path)
        if ext in {".docx", ".doc"}:
            return self._parse_docx(file_path)
        if ext in {".pptx", ".ppt"}:
            return self._parse_pptx(file_path)
        return self._parse_text(file_path)

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    def _init_docling(self) -> bool:
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.document_converter import DocumentConverter, PdfFormatOption

            opts = PdfPipelineOptions(
                do_ocr=False,
                do_table_structure=True,
                generate_page_images=False,
            )
            self._docling = DocumentConverter(
                format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=opts)}
            )
            return True
        except Exception as exc:  # pragma: no cover
            logger.warning("Docling not available, falling back to PyMuPDF: %s", exc)
            self._docling = None
            return False

    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        if self._docling_available and self._docling is not None:
            try:
                return self._parse_pdf_docling(file_path)
            except Exception as exc:
                logger.warning("Docling parse failed, falling back to PyMuPDF: %s", exc)
        return self._parse_pdf_pymupdf(file_path)

    def _parse_pdf_docling(self, file_path: str) -> Dict[str, Any]:
        assert self._docling is not None
        result = self._docling.convert(file_path)
        doc = result.document
        num_pages = len(doc.pages) if doc.pages else 1

        # Bucket text by page_no so we can emit one ``pages`` entry per page.
        pages_content: Dict[int, List[str]] = {i: [] for i in range(1, num_pages + 1)}
        for item in doc.texts:
            for prov in item.prov:
                if prov.page_no in pages_content:
                    pages_content[prov.page_no].append(item.text)

        pages: List[Dict[str, Any]] = []
        for page_num in range(1, num_pages + 1):
            pages.append(
                {
                    "page_num": page_num,
                    "text": "\n\n".join(pages_content[page_num]),
                    "tables": [],
                    "image_refs": [],
                }
            )

        return {
            "text_markdown": doc.export_to_markdown(),
            "pages": pages,
            "metadata": {
                "title": Path(file_path).stem,
                "num_pages": num_pages,
                "file_type": "pdf",
                "source_path": file_path,
            },
        }

    def _parse_pdf_pymupdf(self, file_path: str) -> Dict[str, Any]:
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        pages: List[Dict[str, Any]] = []
        text_content_parts: List[str] = []

        for idx, page in enumerate(doc):
            p_text = page.get_text()
            tables = self._extract_tables_pymupdf(page)
            text_content_parts.append(f"\n--- Page {idx+1} ---\n{p_text}")
            pages.append(
                {
                    "page_num": idx + 1,
                    "text": p_text,
                    "tables": tables,
                    "image_refs": [],
                }
            )

        return {
            "text_markdown": "\n".join(text_content_parts),
            "pages": pages,
            "metadata": {
                "title": Path(file_path).stem,
                "num_pages": len(pages),
                "file_type": "pdf",
                "source_path": file_path,
            },
        }

    @staticmethod
    def _extract_tables_pymupdf(page) -> List[Dict[str, str]]:
        """Lightweight table extraction from a PyMuPDF page.

        Tries ``page.find_tables()``; falls back to empty list if the
        heuristic fails (PyMuPDF's table finder is conservative).
        """
        out: List[Dict[str, str]] = []
        try:
            tables = page.find_tables()
            for idx, t in enumerate(tables):
                df = t.to_pandas() if hasattr(t, "to_pandas") else None
                if df is None:
                    continue
                md = df.to_markdown(index=False)
                out.append({"markdown": md, "caption": f"Detected table {idx + 1}"})
        except Exception:
            return out
        return out

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        from docx import Document as DocxDocument  # python-docx

        document = DocxDocument(file_path)

        # Group paragraphs and tables into pseudo-pages by walking the
        # document body in order. python-docx doesn't expose "pages" but
        # we can still emit one logical page every N paragraphs (or just
        # one page containing the whole document).
        body_paras: List[str] = []
        tables: List[Dict[str, str]] = []
        for element in document.element.body.iter():
            tag = element.tag.split("}", 1)[-1]
            if tag == "p":
                text = "".join(t.text or "" for t in element.iter() if t.tag.endswith("}t"))
                if text.strip():
                    body_paras.append(text)
            elif tag == "tbl":
                # Best-effort: collect the cell texts
                rows = []
                for row in element.iter():
                    if row.tag.endswith("}tr"):
                        cells = [
                            "".join(t.text or "" for t in c.iter() if t.tag.endswith("}t"))
                            for c in row.iter() if c.tag.endswith("}tc")
                        ]
                        if cells:
                            rows.append(cells)
                if rows:
                    md_lines = ["| " + " | ".join(rows[0]) + " |"]
                    md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
                    for row in rows[1:]:
                        md_lines.append("| " + " | ".join(row) + " |")
                    tables.append({"markdown": "\n".join(md_lines), "caption": "DOCX table"})

        full_text = "\n\n".join(body_paras)
        # Heuristic: split on hard page breaks
        page_texts = full_text.split("\x0c")
        if len(page_texts) == 1:
            page_texts = [full_text]

        pages = [
            {"page_num": i + 1, "text": t, "tables": [], "image_refs": []}
            for i, t in enumerate(page_texts)
            if t.strip()
        ]
        # Attach tables to the first page for simplicity
        if pages and tables:
            pages[0]["tables"] = tables
        elif not pages:
            pages = [{"page_num": 1, "text": full_text, "tables": tables, "image_refs": []}]

        return {
            "text_markdown": full_text,
            "pages": pages,
            "metadata": {
                "title": Path(file_path).stem,
                "num_pages": len(pages),
                "file_type": "docx",
                "source_path": file_path,
            },
        }

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        from pptx import Presentation  # python-pptx

        presentation = Presentation(file_path)
        pages: List[Dict[str, Any]] = []
        all_text_parts: List[str] = []

        for idx, slide in enumerate(presentation.slides, start=1):
            text_parts: List[str] = []
            tables: List[Dict[str, str]] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            text_parts.append(line)
                if shape.has_table:
                    tbl = shape.table
                    rows = []
                    for row in tbl.rows:
                        rows.append([cell.text for cell in row.cells])
                    if rows:
                        md = ["| " + " | ".join(rows[0]) + " |"]
                        md.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
                        for r in rows[1:]:
                            md.append("| " + " | ".join(r) + " |")
                        tables.append({"markdown": "\n".join(md), "caption": f"Slide {idx} table"})

            slide_text = "\n".join(text_parts)
            all_text_parts.append(f"\n--- Slide {idx} ---\n{slide_text}")
            pages.append(
                {
                    "page_num": idx,
                    "text": slide_text,
                    "tables": tables,
                    "image_refs": [],
                }
            )

        return {
            "text_markdown": "\n".join(all_text_parts),
            "pages": pages,
            "metadata": {
                "title": Path(file_path).stem,
                "num_pages": len(pages),
                "file_type": "pptx",
                "source_path": file_path,
            },
        }

    # ------------------------------------------------------------------
    # Plain text / markdown
    # ------------------------------------------------------------------

    def _parse_text(self, file_path: str) -> Dict[str, Any]:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        except Exception as exc:
            logger.error("Failed to read %s: %s", file_path, exc)
            text = f"Error reading file: {exc}"

        ext = Path(file_path).suffix.lower().lstrip(".")
        pages = [{"page_num": 1, "text": text, "tables": [], "image_refs": []}]

        return {
            "text_markdown": text,
            "pages": pages,
            "metadata": {
                "title": Path(file_path).stem,
                "num_pages": 1,
                "file_type": ext,
                "source_path": file_path,
            },
        }
